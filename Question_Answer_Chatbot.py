import os
import tempfile
import streamlit as st
import google.generativeai as genai
import fitz
import docx
from pptx import Presentation
from io import BytesIO
import requests
import re

# Set environment variables for Google API
os.environ["GOOGLE_API_KEY"] = "AIzaSyAWjOyvXsq6oq_uhduhvP1i4sbYEmBgN1I"
os.environ["GOOGLE_CSE_ID"] = "AIzaSyAWjOyvXsq6oq_uhduhvP1i4sbYEmBgN1I"

google_api_key = os.getenv("GOOGLE_API_KEY")
google_cse_id = os.getenv("GOOGLE_CSE_ID")

if not google_api_key or not google_cse_id:
    raise ValueError("GOOGLE_API_KEY or GOOGLE_CSE_ID not found in environment variables.")
genai.configure(api_key=google_api_key)

def perform_web_search(query):
    search_url = "https://www.googleapis.com/customsearch/v1"
    params = {
        "q": query,
        "cx": google_cse_id,
        "key": google_api_key,
        "num": 3,
    }
    response = requests.get(search_url, params=params)
    results = response.json()
    search_results = ""
    for item in results.get("items", []):
        search_results += item["snippet"] + "\n"
    return search_results

def extract_text_from_pdf(file):
    text = ""
    pdf_document = fitz.open(stream=file.read(), filetype="pdf")
    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        text += page.get_text()
    return text

def extract_text_from_docx(file):
    doc = docx.Document(file)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

def extract_text_from_pptx(file):
    text = ""
    try:
        presentation = Presentation(file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        text = f"Failed to extract text from PPTX: {e}"
    return text

def extract_text_from_txt(file):
    text = file.read().decode("utf-8")
    return text

file_type_handlers = {
    "pdf": extract_text_from_pdf,
    "docx": extract_text_from_docx,
    "pptx": extract_text_from_pptx,
    "txt": extract_text_from_txt,
}

def extract_text(file, file_type):
    handler = file_type_handlers.get(file_type)
    if handler:
        return handler(file)
    else:
        return "Unsupported file type. Please try a PDF, DOCX, PPTX, or TXT file."

def summarize_text(text):
    model = genai.GenerativeModel('gemini-pro')
    response = model.generate_content(f"Summarize the following text: {text}")
    return response.text

def summarize_based_on_topics(text, topics):
    model = genai.GenerativeModel('gemini-pro')
    response = model.generate_content(f"Summarize the following text focusing on the topics {topics}: {text}")
    return response.text

def explain_concept(concept, text):
    web_search_results = perform_web_search(concept)
    combined_text = f"Based on the document:\n{text}\n\nAnd additional information from the web:\n{web_search_results}"
    model = genai.GenerativeModel('gemini-pro')
    response = model.generate_content(f"Explain the concept of {concept} based on the following information: {combined_text}")
    return response.text

def get_gemini_response(question, text):
    model = genai.GenerativeModel('gemini-pro')
    response = model.generate_content(f"Answer the following question based on the document: {question}. Document: {text}")
    return response.text

def generate_custom_quiz(topic, text):
    model = genai.GenerativeModel('gemini-pro')
    response = model.generate_content(f"Generate quiz questions based on the topic '{topic}' from the following text: {text}")
    questions = response.text.strip().split("\n")
    cleaned_questions = []
    for question in questions:
        if "?" in question:
            question_clean = re.sub(r"^\d+\.\s*|\d+\s*\.\s*", "", question)
            question_clean = re.sub(r"^Question\s*\d+\s*:\s*", "", question_clean)
            cleaned_questions.append(question_clean.strip())

    return cleaned_questions[:10] 

st.set_page_config(page_title="Study Helper")
st.header("Study Helper")
uploaded_file = st.file_uploader("Upload your document", type=list(file_type_handlers.keys()))

if uploaded_file:
    file_type = uploaded_file.name.split('.')[-1].lower()
    text = extract_text(uploaded_file, file_type)
    
    if "Failed" not in text:
        summary = summarize_text(text)
        st.subheader("Document Summary")
        st.write(summary)

        topics = st.text_input("Enter topics for focused summarization (comma-separated):", key="topics")
        if st.button("Summarize Based on Topics"):
            if topics:
                topic_summary = summarize_based_on_topics(text, topics)
                st.subheader("Topic-Based Summary")
                st.write(topic_summary)

        concept = st.text_input("Enter a concept to get an explanation:", key="concept")
        if st.button("Explain Concept"):
            if concept:
                explanation = explain_concept(concept, text)
                st.subheader("Concept Explanation")
                st.write(explanation)

        custom_topic = st.text_input("Enter a topic for custom quiz generation:", key="custom_topic")

        if st.button("Generate Custom Quiz"):
            if custom_topic:
                quiz_questions = generate_custom_quiz(custom_topic, text)
                st.session_state.quiz_questions = quiz_questions  
                st.session_state.user_answers = [""] * len(quiz_questions)  
        if "quiz_questions" in st.session_state:
            st.subheader("Custom Quiz Questions")
            for i, question in enumerate(st.session_state.quiz_questions):
                st.write(f"**Question {i+1}:** {question}")
                st.session_state.user_answers[i] = st.text_input(f"Your answer to question {i+1}:", key=f"answer_{i}", value=st.session_state.user_answers[i])

            if st.button("Submit Quiz"):
                correct_answers = len([answer for answer in st.session_state.user_answers if answer.strip()]) 
                total_questions = len(st.session_state.quiz_questions)
                st.write(f"You answered {correct_answers} out of {total_questions} questions.")

        question = st.text_input("Ask a question based on the document:", key="question_ask")
        if question:
            with st.spinner("Getting response..."):
                response = get_gemini_response(question, text)
                st.subheader("Answer")
                st.write(response)
    else:
        st.error(text)
